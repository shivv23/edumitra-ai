from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── colors ──
BG_DARK = RGBColor(0x0B, 0x0B, 0x1A)
BG_CARD = RGBColor(0x12, 0x12, 0x2E)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)  # cyan
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)  # purple
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xCC)
SUBTLE = RGBColor(0x66, 0x66, 0x99)

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, title, bullets, subtext=None):
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), title,
                 font_size=36, color=ACCENT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT2
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
        p.space_after = Pt(12)
    if subtext:
        p2 = tf.add_paragraph()
        p2.text = subtext
        p2.font.size = Pt(14)
        p2.font.color.rgb = SUBTLE
        p2.space_before = Pt(20)

def add_card_grid(slide, items, start_top=1.6):
    """items: [(title, desc), ...] in 2x2 grid"""
    for i, (title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 6.2)
        top = Inches(start_top + row * 2.6)
        card = add_shape(slide, left, top, Inches(5.8), Inches(2.2), BG_CARD)
        add_text_box(slide, left + Inches(0.3), top + Inches(0.2), Inches(5.2), Inches(0.5),
                     title, font_size=22, color=ACCENT, bold=True)
        add_text_box(slide, left + Inches(0.3), top + Inches(0.8), Inches(5.2), Inches(1.2),
                     desc, font_size=15, color=GRAY)

# ══════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Pt(6))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "EduMitra AI", font_size=60, color=WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
             "Personalized Learning & Mental Wellness Companion\nfor Indian Students",
             font_size=26, color=ACCENT)

card = add_shape(slide, Inches(1), Inches(4.2), Inches(7), Inches(1.5), BG_CARD)
add_text_box(slide, Inches(1.3), Inches(4.4), Inches(6.5), Inches(1.2),
             "Team Name: shivamkumar0423\nTeam Leader: Shivam Kumar\nProposal Submission",
             font_size=18, color=GRAY)

shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.5), Inches(1.5), Pt(6), Inches(4.5))
shape2.fill.solid()
shape2.fill.fore_color.rgb = ACCENT2
shape2.line.fill.background()

# ══════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Problem Statement", [
    "▸ 250M+ Indian students across fragmented curricula (CBSE, state boards, competitive exams) with no personalized learning adaptation",
    "▸ 53% of Indian students report anxiety — yet wellness support is completely siloed from academics",
    "▸ Existing edtech platforms offer one-size-fits-all content: no adaptation to individual pace or knowledge gaps",
    "▸ No single platform understands both a student's learning needs and emotional state simultaneously",
    "▸ DPDP Act 2023 compliance is an afterthought — student data privacy is at risk in most edtech solutions",
])

# ══════════════════════════════════════════
# SLIDE 3 — PROPOSED SOLUTION
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Proposed Solution", font_size=36, color=ACCENT, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT2
shape.line.fill.background()

add_card_grid(slide, [
    ("Learning Agent", "RAG-augmented LLM (Claude + Gemini) generates\npersonalized study plans, quizzes, and flashcards\nfrom NCERT content via ChromaDB vector store"),
    ("Wellness Agent", "Rule-based crisis detection + sentiment analysis.\nAuto-alerts at medium/high risk. Encrypted at rest,\ncompletely excluded from logs."),
    ("Voice Agent", "11 Indian languages via Web Speech API +\nSarvam AI STT fallback. Enables access for\nvernacular-medium and low-literacy students."),
    ("Data Protection", "DPDP Act 2023: right to erasure, data portability,\nencryption at rest, parental consent gate\nfor minors under 18."),
])

# ══════════════════════════════════════════
# SLIDE 4 — TECHNOLOGY STACK
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Technology Stack", font_size=36, color=ACCENT, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT2
shape.line.fill.background()

add_card_grid(slide, [
    ("Frontend", "Next.js 14 (App Router) · TypeScript\nMaterial UI · Glassmorphism design\nHosted on Vercel (global CDN)"),
    ("Backend", "FastAPI (Python 3.11)\nLangGraph (multi-agent orchestration)\nHosted on Railway (Docker/Nixpacks)"),
    ("AI / ML", "Claude Sonnet 4 (primary LLM)\nGemini 2.5 Flash (fallback LLM)\nChromaDB vector store (all-MiniLM-L6-v2)"),
    ("APIs & Storage", "Sarvam AI (STT) · Stability AI (image gen)\nSupabase PostgreSQL (managed DB + RLS)\nWeb Speech API (browser-side STT)"),
])

# ══════════════════════════════════════════
# SLIDE 5 — FINAL DELIVERABLES
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Final Deliverables", font_size=36, color=ACCENT, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT2
shape.line.fill.background()

deliverables = [
    ("Working Web App", "Fully functional Next.js frontend deployed on Vercel\nwith glassmorphism UI, real-time study plans,\nwellness check-ins, and teacher/parent dashboards"),
    ("Multi-Agent Backend", "FastAPI server with 4 LangGraph agents running\non Railway — Learning, Wellness, Voice, and\nData Protection — all operational end-to-end"),
    ("RAG Knowledge Base", "ChromaDB vector store seeded with NCERT content\nacross 5+ subjects. Semantic retrieval with\n0.65 relevance threshold for accurate answers"),
    ("Wellness Alert System", "Real-time sentiment analysis with auto-escalation.\nMedium/high risk triggers alerts. All wellness data\nencrypted at rest with strictest access controls."),
    ("Voice Interface", "Browser-based STT supporting 11 Indian languages\nwith Sarvam AI backend fallback. Enables\nvernacular-medium students to interact naturally"),
    ("DPDP Act Compliance", "Right to erasure (DELETE /api/data/me), data\nportability (GET /api/data/export), encryption\nat rest, parental consent gate for minors"),
]

for i, (title, desc) in enumerate(deliverables):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.15)
    top = Inches(1.5 + row * 2.8)
    card = add_shape(slide, left, top, Inches(3.9), Inches(2.5), BG_CARD)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(3.5), Inches(0.5),
                 title, font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.8), Inches(3.5), Inches(1.5),
                 desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 6 — EXPECTED IMPACT
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Expected Impact", [
    "▸ Personalized learning at scale — 250M+ students get adaptive study plans tailored to their pace and knowledge gaps",
    "▸ Early wellness intervention — burnout risk detected before academic performance drops, with immediate alerts",
    "▸ Bridging the language gap — vernacular voice interface enables access for non-English-medium students across India",
    "▸ Teacher empowerment — real-time class-wide progress heatmaps and at-risk student identification in one dashboard",
    "▸ Parent visibility — academic + wellness insights without breaching the child's data privacy",
    "▸ Privacy-first architecture — full DPDP Act 2023 compliance built in from Day 1, not retrofitted as an afterthought",
])

# ══════════════════════════════════════════
# SLIDE 7 — PROJECT ROADMAP
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Project Roadmap", font_size=36, color=ACCENT, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT2
shape.line.fill.background()

phases = [
    ("Foundation\n(Weeks 1–2)", "• Monorepo setup & CI/CD\n• Learning Agent: RAG pipeline\n  + study plan generation\n• Quiz engine with adaptive\n  difficulty\n• ChromaDB seeding"),
    ("Core Features\n(Weeks 3–4)", "• Wellness Agent: sentiment\n  analysis & crisis detection\n• Voice Agent: Web Speech API\n  + Sarvam STT + Gemini chat\n• Teacher dashboard with\n  progress heatmaps"),
    ("Privacy & Polish\n(Weeks 5–6)", "• Data Protection Module:\n  erasure, portability, consent\n• Parent dashboard\n• End-to-end testing\n• Bug fixes & performance\n  optimization\n• Demo preparation"),
]

for i, (phase, desc) in enumerate(phases):
    left = Inches(0.8 + i * 4.15)
    card = add_shape(slide, left, Inches(1.5), Inches(3.9), Inches(5.2), BG_CARD)
    add_text_box(slide, left + Inches(0.2), Inches(1.7), Inches(3.5), Inches(0.8),
                 phase, font_size=20, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), Inches(2.6), Inches(3.5), Inches(3.8),
                 desc, font_size=15, color=GRAY, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# SLIDE 8 — UNIQUENESS & INNOVATION
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "What Makes EduMitra Unique?", [
    "▸ First platform to unify academic personalization + mental wellness in a single AI-driven system",
    "▸ Multi-agent architecture (LangGraph) — specialized agents collaborate, not one monolithic model",
    "▸ RAG over NCERT curriculum ensures answers are syllabus-aligned, not generic web scrapings",
    "▸ Voice-first design in 11 Indian languages — built for Bharat, not just English-speaking users",
    "▸ Privacy-by-design: DPDP Act 2023 compliance built into the architecture from Day 1",
    "▸ Fully self-contained demo mode — presentation works offline without any backend dependency",
])

# ══════════════════════════════════════════
# SLIDE 9 — THANK YOU
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), prs.slide_width, Pt(6))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.2),
             "Thank You", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.6),
             "Team shivamkumar0423  |  Shivam Kumar",
             font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
             "EduMitra AI — Personalized Learning & Wellness for Every Indian Student",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# ── save ──
out = "C:\\Users\\HP\\Desktop\\HackArena\\edumitra-ai\\docs\\week0\\EduMitra_AI_Proposal.pptx"
prs.save(out)
print(f"Saved: {out}")
